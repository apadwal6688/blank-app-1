import streamlit as st
import pptx
from pptx import Presentation
from gtts import gTTS
import os
import tempfile
from pathlib import Path
import uuid
import anthropic

class ClaudePPTScriptGenerator:
    def __init__(self, api_key):
        """
        Initialize the Claude-powered script generator
        
        Args:
            api_key (str): Anthropic Claude API key
        """
        self.client = anthropic.Anthropic(api_key=api_key)
        self.model = "claude-3-5-sonnet-20240620"
    
    def generate_slide_script(self, slide_text):
        """
        Generate a professional presentation script for a slide using Claude AI
        
        Args:
            slide_text (str): Text content of the slide
        
        Returns:
            str: Generated script for the slide
        """
        try:
            response = self.client.messages.create(
                model=self.model,
                max_tokens=300,
                messages=[
                    {
                        "role": "user",
                        "content": f"""Create a professional, engaging presentation script for a slide with the following content:

Slide Content: {slide_text}

Requirements:
- Write in a clear, confident speaking style
- Provide context and explain key points
- Use professional language suitable for business or academic presentations
- Include transitions and highlights
- Aim for about 2-3 sentences per slide"""
                    }
                ]
            )
            return response.content[0].text
        except Exception as e:
            st.error(f"Error generating script with Claude: {e}")
            return f"Slide content: {slide_text}"

def extract_slide_text(presentation):
    """
    Extract text from each slide in the PowerPoint presentation.
    
    Args:
        presentation (Presentation): Loaded PowerPoint presentation
    
    Returns:
        list: List of text content for each slide
    """
    slide_texts = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slide_texts.append(" ".join(slide_text))
    return slide_texts

def generate_audio(scripts, voice='en'):
    """
    Generate audio files for each slide script.
    
    Args:
        scripts (list): List of slide scripts
        voice (str): Language code for text-to-speech
    
    Returns:
        list: Paths to generated audio files
    """
    audio_files = []
    for i, script in enumerate(scripts, 1):
        # Create a unique filename
        filename = f"slide_{i}_{uuid.uuid4()}.mp3"
        temp_audio_path = os.path.join(tempfile.gettempdir(), filename)
        
        # Use gTTS to generate audio
        tts = gTTS(text=script, lang=voice)
        tts.save(temp_audio_path)
        audio_files.append(temp_audio_path)
    
    return audio_files

def main():
    st.set_page_config(page_title="Claude PPT Script Generator", page_icon="🎤")
    
    st.title("Claude AI PowerPoint Script Generator 🎤📊")
    
    # API Key Input
    api_key = st.text_input("Enter your Anthropic Claude API Key", type="password")
    
    # File upload and processing
    uploaded_file = st.file_uploader("Upload PowerPoint Presentation", type=['pptx'])
    
    # Voice selection
    voice_options = {
        'English (US)': 'en',
        'Spanish': 'es',
        'French': 'fr',
        'German': 'de',
        'Italian': 'it',
        'Portuguese': 'pt'
    }
    selected_voice = st.selectbox("Select Voice Language", list(voice_options.keys()))
    
    # Process button
    process_button = st.button("Generate Scripts")
    
    if process_button and api_key and uploaded_file is not None:
        # Initialize Claude Script Generator
        script_generator = ClaudePPTScriptGenerator(api_key)
        
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name
        
        try:
            # Load presentation
            presentation = Presentation(tmp_file_path)
            
            # Extract slide texts
            slide_texts = extract_slide_text(presentation)
            
            # Generate scripts using Claude
            scripts = []
            progress_bar = st.progress(0)
            for i, text in enumerate(slide_texts):
                script = script_generator.generate_slide_script(text)
                scripts.append(script)
                progress_bar.progress((i + 1) / len(slide_texts))
            
            # Generate audio
            audio_files = generate_audio(scripts, voice=voice_options[selected_voice])
            
            # Display results
            st.header("📜 Generated Presentation Scripts")
            for i, script in enumerate(scripts, 1):
                st.subheader(f"Slide {i}")
                st.write(script)
                
                # Play audio
                st.audio(audio_files[i-1], format='audio/mp3')
            
        except Exception as e:
            st.error(f"An error occurred: {e}")
        
        finally:
            # Cleanup temporary files
            if 'tmp_file_path' in locals():
                os.unlink(tmp_file_path)
            if 'audio_files' in locals():
                for audio_file in audio_files:
                    os.unlink(audio_file)

if __name__ == "__main__":
    main()